import streamlit as st
from pptx import Presentation
from docx import Document
import tempfile
from transformers import pipeline

# ——— Load the HF generator once at startup ———
@st.cache_resource(show_spinner=False)
def get_generator():
    return pipeline("text2text-generation", model="google/flan-t5-small")

generator = get_generator()

st.set_page_config(page_title="Proposal Slide Generator", layout="wide")
st.title("Proposal Slide Text Generator")

# ——— Sidebar file uploads ———
st.sidebar.header("Upload Files")
pptx_files = st.sidebar.file_uploader(
    "Upload past PPTX slides", type=["pptx"], accept_multiple_files=True
)
brief_file = st.sidebar.file_uploader(
    "Upload project brief (DOCX or TXT)", type=["docx", "txt"]
)

# ——— Utility: extract Objectives & Solutions text from PPTX ———
def extract_text_from_pptx(pptx_streams):
    examples = []
    for f in pptx_streams:
        prs = Presentation(f)
        for slide in prs.slides:
            text_bits = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_bits.append(shape.text)
            slide_text = "\n".join(text_bits)
            if "Objective" in slide_text or "Solution" in slide_text:
                examples.append(slide_text)
    return examples

# ——— Utility: extract text from DOCX or TXT ———
def extract_text_from_brief(brief_stream, filename):
    if filename.lower().endswith(".docx"):
        doc = Document(brief_stream)
        return "\n".join([p.text for p in doc.paragraphs])
    else:
        return brief_stream.read().decode("utf-8")

# ——— Main generation flow ———
if st.sidebar.button("Generate Slide Text"):
    if not pptx_files or not brief_file:
        st.sidebar.error("Please upload both past PPTX files and a project brief.")
    else:
        with st.spinner("Extracting examples..."):
            examples = extract_text_from_pptx(pptx_files)
            brief_text = extract_text_from_brief(brief_file, brief_file.name)

        prompt = (
            "You are a professional proposal-writer assistant. Based on the following "
            "Objectives & Solutions examples and the new project brief, draft a concise "
            "Objectives & Solutions paragraph in our brand tone.\n\n"
            "=== Examples ===\n"
            + "\n---\n".join(examples)
            + "\n\n=== New Brief ===\n"
            + brief_text
            + "\n\n=== Draft ===\n"
        )

        with st.spinner("Generating draft..."):
            result = generator(prompt, max_length=256)[0]
            draft = result["generated_text"].strip()

        st.subheader("Draft Objectives & Solutions Text")
        draft_edited = st.text_area("Edit as needed:", draft, height=300)

        # ——— Download buttons ———
        col1, col2 = st.columns(2)
        with col1:
            if st.button("Download as DOCX"):
                doc = Document()
                doc.add_paragraph(draft_edited)
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
                doc.save(tmp.name)
                st.download_button(
                    "Download DOCX",
                    data=open(tmp.name, "rb").read(),
                    file_name="Objectives_Solutions.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
        with col2:
            if st.button("Download as TXT"):
                tmp_txt = tempfile.NamedTemporaryFile(delete=False, suffix=".txt")
                tmp_txt.write(draft_edited.encode("utf-8"))
                tmp_txt.close()
                st.download_button(
                    "Download TXT",
                    data=open(tmp_txt.name, "rb").read(),
                    file_name="Objectives_Solutions.txt",
                    mime="text/plain",
                )

